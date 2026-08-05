"""PPTX 解析器 — 基于 python-pptx。

逐 slide 提取文本框内容,每个 slide 的文本作为一个或多个 chunk。
页码 = slide 序号(1-based),行号在 slide 内从 1 开始。
"""

from __future__ import annotations

import logging
from pathlib import Path

from app.fileparse.base import BaseParser, ParseChunk

logger = logging.getLogger("lqd.fileparse.pptx")


class PptxParser(BaseParser):
    extensions = [".pptx"]

    def parse(self, file_path: Path) -> list[ParseChunk]:
        from pptx import Presentation

        logger.debug("PPTX 解析开始: %s", file_path.name)
        chunks: list[ParseChunk] = []
        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            logger.error("PPTX 打开失败 %s: %s", file_path, exc)
            raise

        slide_count = len(prs.slides)
        logger.info("PPTX %s: 共 %d 张幻灯片", file_path.name, slide_count)

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            shape_count = 0

            for shape in slide.shapes:
                shape_count += 1
                # 文本框
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                # 表格
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))

            logger.debug(
                "PPTX %s slide %d: %d shapes, %d text lines",
                file_path.name, slide_idx, shape_count, len(slide_texts)
            )

            if not slide_texts:
                logger.debug("PPTX %s slide %d: 无文本内容(可能是图片/图表)", file_path.name, slide_idx)
                continue

            # slide 内文本按目标字符数切片
            slide_chunks = self._split_text_to_chunks(
                text="\n".join(slide_texts),
                page=slide_idx,
                page_label=f"幻灯片 {slide_idx}",
                line_offset=1,
            )
            chunks.extend(slide_chunks)

        logger.info("PPTX 解析完成: %s, %d/%d slides 有文本, %d 切片",
                    file_path.name, len({c.page for c in chunks}), slide_count, len(chunks))
        return chunks
