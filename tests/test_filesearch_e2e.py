"""端到端测试：创建测试文件 → 构建索引 → 检索 → 验证结果"""
import sys
import os
import shutil
import json
from pathlib import Path

# 确保项目根在 sys.path
project_root = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(project_root))

from docx import Document
from pptx import Presentation
from openpyxl import Workbook

from app.fileparse import parse_file, SUPPORTED_EXTENSIONS
from app.fileindex.builder import build_full, build_incremental
from app.fileindex.store import FileIndexStore
from app.fileindex.status import start_build, get_status

TEST_DIR = project_root / "data" / "test_files"
INDEX_DIR = project_root / "data" / "fileindex"

def create_test_files():
    """创建测试文件"""
    if TEST_DIR.exists():
        shutil.rmtree(TEST_DIR)
    TEST_DIR.mkdir(parents=True, exist_ok=True)

    # DOCX
    doc = Document()
    doc.add_paragraph("这是第一个段落，包含关键词 iSC 平台架构")
    doc.add_paragraph("第二个段落讨论门禁系统配置")
    doc.add_paragraph("最后段落涉及视频监控和报警联动")
    doc.save(str(TEST_DIR / "test_doc.docx"))

    # PPTX
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "iSC 平台概述"
    slide1.placeholders[1].text = "本页介绍 iSC 综合安防管理平台的核心架构"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "门禁管理"
    slide2.placeholders[1].text = "门禁系统配置与访客管理功能"
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "视频监控"
    slide3.placeholders[1].text = "视频监控与报警联动机制"
    prs.save(str(TEST_DIR / "test_ppt.pptx"))

    # XLSX
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "设备清单"
    ws1.append(["设备名称", "类型", "位置"])
    ws1.append(["摄像头A", "视频", "大门口"])
    ws1.append(["门禁控制器B", "门禁", "前台"])
    ws2 = wb.create_sheet("报警记录")
    ws2.append(["时间", "事件", "处理人"])
    ws2.append(["2024-01-01", "报警触发", "管理员"])
    wb.save(str(TEST_DIR / "test_excel.xlsx"))

    # TXT
    with open(str(TEST_DIR / "test_notes.txt"), "w", encoding="utf-8") as f:
        f.write("iSC 平台部署文档\n门禁系统配置流程\n视频监控参数调优\n报警联动测试结果\n")

    print(f"[OK] 测试文件创建完成: {TEST_DIR}")
    for p in sorted(TEST_DIR.iterdir()):
        print(f"  - {p.name} ({p.stat().st_size} bytes)")

def test_parse():
    """测试文件解析"""
    print("\n=== 测试文件解析 ===")
    for ext in [".docx", ".pptx", ".xlsx", ".txt"]:
        files = list(TEST_DIR.glob(f"*{ext}"))
        if not files:
            continue
        f = files[0]
        result = parse_file(f)
        if result.error:
            print(f"[FAIL] {f.name}: {result.error}")
            continue
        print(f"[OK] {f.name}: {len(result.chunks)} 切片")
        for i, c in enumerate(result.chunks):
            print(f"  chunk {i}: page={c.page} label='{c.page_label}' "
                  f"lines={c.line_start}-{c.line_end} text='{c.text[:50]}...'")

def test_build_full():
    """测试全量索引构建"""
    print("\n=== 测试全量索引构建 ===")
    if INDEX_DIR.exists():
        shutil.rmtree(INDEX_DIR)
    INDEX_DIR.mkdir(parents=True, exist_ok=True)

    result = build_full(str(TEST_DIR), str(INDEX_DIR))
    print(f"ok={result.ok} files_scanned={result.files_scanned} "
          f"files_indexed={result.files_indexed} chunks_built={result.chunks_built} "
          f"duration={result.duration_sec:.2f}s")
    if result.error:
        print(f"[ERROR] {result.error}")
    print("--- 构建日志 ---")
    for line in result.log:
        print(f"  {line}")
    return result

def test_search():
    """测试检索"""
    print("\n=== 测试检索 ===")
    store = FileIndexStore(str(INDEX_DIR))
    data = store.load()
    print(f"索引概况: files={len(data.files)} chunks={len(data.chunks)} "
          f"tokens={len(data.inverted_index)}")

    # 模拟关键词检索
    from app.retrieval.tokenizer import tokenize_unique, expand_query_weighted
    query = "门禁系统"
    tokens = tokenize_unique(query)
    expanded_tokens, weights = expand_query_weighted(tokens, query)
    print(f"查询: '{query}' → tokens={tokens} expanded_tokens={expanded_tokens} weights={weights}")

    # 收集候选 chunk_id
    candidate_ids = set()
    for term in expanded_tokens:
        chunk_ids = data.inverted_index.get(term, [])
        for cid in chunk_ids:
            candidate_ids.add(cid)
    print(f"候选切片数: {len(candidate_ids)}")

    # 简单打分
    chunk_map = {c.chunk_id: c for c in data.chunks}
    results = []
    for cid in candidate_ids:
        c = chunk_map.get(cid)
        if c:
            score = sum(w for term, w in weights.items()
                       if term in c.text.lower())
            results.append((score, c))
    results.sort(key=lambda x: -x[0])

    print(f"\n检索结果 ({len(results)} 条):")
    for rank, (score, c) in enumerate(results[:10], 1):
        print(f"  #{rank} score={score:.2f} file={c.file_name} "
              f"page={c.page}({c.page_label}) lines={c.line_start}-{c.line_end}")
        print(f"       text: {c.text[:80]}...")

def test_incremental():
    """测试增量索引"""
    print("\n=== 测试增量索引(无变更) ===")
    result = build_incremental(str(TEST_DIR), str(INDEX_DIR))
    print(f"ok={result.ok} scanned={result.files_scanned} "
          f"indexed={result.files_indexed} skipped={result.files_skipped} "
          f"chunks={result.chunks_built} duration={result.duration_sec:.2f}s")

def test_status_api():
    """测试状态管理 API"""
    print("\n=== 测试状态管理 API ===")
    job_id = start_build("incremental", str(TEST_DIR), str(INDEX_DIR))
    print(f"启动任务: job_id={job_id}")
    import time
    time.sleep(1)
    status = get_status(job_id)
    if status:
        print(f"状态: {status['status']} stage={status.get('current_stage','')}")
        if status.get('result'):
            r = status['result']
            print(f"  结果: ok={r['ok']} files={r['files_indexed']} "
                  f"chunks={r['chunks_built']}")

if __name__ == "__main__":
    create_test_files()
    test_parse()
    test_build_full()
    test_search()
    test_incremental()
    test_status_api()
    print("\n=== 端到端测试完成 ===")
